# -*- coding: utf-8 -*-
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # 建立簡報並設定為 16:9 寬螢幕尺寸
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 常用配色
    BG_COLOR = RGBColor(20, 22, 28)        # #14161C (極深暗藍黑色)
    CARD_BG_COLOR = RGBColor(30, 34, 42)   # #1E222A (卡片深灰色)
    BORDER_COLOR = RGBColor(45, 50, 60)    # 卡片預設邊框色
    
    CYBER_BLUE = RGBColor(0, 180, 255)     # #00B4FF (科技藍)
    CYBER_RED = RGBColor(255, 60, 60)      # #FF3C3C (警報紅)
    CYBER_GREEN = RGBColor(0, 255, 120)    # #00FF78 (安全綠)
    CYBER_GRAY = RGBColor(160, 165, 175)   # #A0A5AF (輔助灰色)
    
    TEXT_LIGHT = RGBColor(240, 240, 245)   # #F0F0F5 (主文字白色)
    TEXT_MUTED = RGBColor(180, 185, 195)   # #B4B9C3 (次文字淡灰)
    
    # ------------------ 輔助函式 ------------------
    
    def add_background(slide):
        """為投影片加上深色背景"""
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background() # 無邊框
        return bg

    def add_header(slide, title_text):
        """為投影片加上科技感標頭"""
        # 頂部裝飾細線
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = CYBER_BLUE
        line.line.fill.background()
        
        # 標題文字方塊
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), prs.slide_width - Inches(1.2), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = CYBER_BLUE

    def draw_card(slide, left, top, width, height, title, body_list, accent_color=CYBER_BLUE):
        """繪製單個卡片區塊"""
        # 卡片底色與邊框
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG_COLOR
        card.line.color.rgb = accent_color
        card.line.width = Pt(1.5)
        
        # 文字方塊
        txBox = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), height - Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
        
        # 卡片標題
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = 'Microsoft JhengHei'
        p_title.font.size = Pt(17)
        p_title.font.bold = True
        p_title.font.color.rgb = accent_color
        
        # 卡片內容
        for text in body_list:
            p = tf.add_paragraph()
            p.text = text
            p.font.name = 'Microsoft JhengHei'
            p.font.size = Pt(11.5)
            p.font.color.rgb = TEXT_LIGHT
            p.space_before = Pt(6)
            
    # ------------------ SLIDE 1: 封面 ------------------
    blank_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_layout)
    add_background(slide1)
    
    # 裝飾性外框
    glow = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), prs.slide_width - Inches(1.0), prs.slide_height - Inches(1.0))
    glow.fill.background()
    glow.line.color.rgb = CYBER_BLUE
    glow.line.width = Pt(2)
    
    # 裝飾性四角裝飾線
    corner_box = slide1.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(3.0), Inches(0.4))
    corner_box.text_frame.text = "[ SURVEILLANCE ACTIVE ]"
    p_c = corner_box.text_frame.paragraphs[0]
    p_c.font.name = 'Courier New'
    p_c.font.size = Pt(10)
    p_c.font.bold = True
    p_c.font.color.rgb = CYBER_GREEN
    
    # 封面主標題與副標題
    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(2.1), prs.slide_width - Inches(2.0), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p_main = tf.paragraphs[0]
    p_main.text = "🛡️ AI Soft-NVR\n智慧安全監控與入侵辨識系統"
    p_main.alignment = PP_ALIGN.CENTER
    p_main.font.name = 'Microsoft JhengHei'
    p_main.font.size = Pt(38)
    p_main.font.bold = True
    p_main.font.color.rgb = CYBER_BLUE
    
    p_space = tf.add_paragraph()
    p_space.text = ""
    p_space.font.size = Pt(14)
    
    p_sub = tf.add_paragraph()
    p_sub.text = "基於 FastAPI, OpenCV 與 YOLOv8 的低延遲、高安全性居家防衛方案"
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.font.name = 'Microsoft JhengHei'
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = TEXT_LIGHT
    
    p_space2 = tf.add_paragraph()
    p_space2.text = ""
    p_space2.font.size = Pt(20)
    
    p_info = tf.add_paragraph()
    p_info.text = "專案核心概念 ・ 修改演進歷程 ・ 關鍵問題與解決方案"
    p_info.alignment = PP_ALIGN.CENTER
    p_info.font.name = 'Microsoft JhengHei'
    p_info.font.size = Pt(13)
    p_info.font.color.rgb = CYBER_GREEN
    
    # ------------------ SLIDE 2: 核心特色 ------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_background(slide2)
    add_header(slide2, "🌟 專案核心概念與特色")
    
    # 左卡片: 即時低延遲影像串流
    draw_card(
        slide2,
        left=Inches(0.6),
        top=Inches(1.5),
        width=Inches(5.8),
        height=Inches(5.2),
        title="⚡ 即時低延遲影像串流",
        body_list=[
            "• WebSockets 傳輸技術：後端擷取影像進行高效 JPEG 壓縮，轉為 Base64 編碼即時推播，FPS 達 ~20 影格，維持超低延遲即時畫面體驗。",
            "• 電子藍（Cyber Blue）即時標記：即時影像動態疊加 AI 藍框與偵測信心度標籤；在觸發事件錄影時，於網頁上加註紅色閃爍「● REC」錄影字樣。",
            "• 異常連線優雅釋放：WebSocket 連線若因瀏覽器關閉/重新整理而中斷，後端能自動捕獲異常並立刻終止背景迴圈，防止 CPU 空轉與 Log 狂飆。"
        ],
        accent_color=CYBER_BLUE
    )
    
    # 右卡片: 雙重 AI 入侵辨識引擎與模擬器
    draw_card(
        slide2,
        left=Inches(6.9),
        top=Inches(1.5),
        width=Inches(5.8),
        height=Inches(5.2),
        title="🧠 雙重 AI 偵測與雷達模擬器",
        body_list=[
            "• 高級 YOLOv8 模式：系統啟動時若偵測到本機有 ultralytics 依賴，會自動加載並執行 yolov8n 奈米模型，進行高精度人體目標鎖定，防範因風吹草動、寵物或光線引起的誤警報。",
            "• 預設輕量模式：無 YOLO 環境時自動切換至內建 OpenCV Haar Cascades 人臉識別與動態偵測，兼顧輕量 CPU（如 Raspberry Pi）執行需求。",
            "• 內建雷達模擬器 (Simulator)：專門為無實體相機鏡頭的開發環境設計，提供高科技感雷達掃描背景與移動入侵點模擬，實現無硬體相機的零門檻敏捷開發。"
        ],
        accent_color=CYBER_GREEN
    )

    # ------------------ SLIDE 3: 兩大獨家核心創新技術 ------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_background(slide3)
    add_header(slide3, "⏪ 兩大獨家核心創新技術")
    
    # 左卡片: 5秒預錄回溯緩衝區
    draw_card(
        slide3,
        left=Inches(0.6),
        top=Inches(1.5),
        width=Inches(5.8),
        height=Inches(5.2),
        title="⏪ 獨家 5 秒「預錄回溯緩衝區」",
        body_list=[
            "• 記憶體環形緩衝快取：在記憶體中維護一個雙端佇列（deque(maxlen=100)），在 20 FPS 下滾動式儲存最近 5 秒的連續監視影格。",
            "• 事件前因還原：當 AI 偵測判定入侵發生並啟動錄影時，系統會「先將 5 秒前記憶體快取的畫面完整寫入影片」，隨後再繼續錄製動態入侵畫面。",
            "• 保障黃金關鍵：這徹底解決了以往「偵測到入侵才錄影、卻已錯過嫌犯如何翻牆或接近的最關鍵 5 秒前置動態」的痛點，完整保留事件前因後果。"
        ],
        accent_color=CYBER_GREEN
    )
    
    # 右卡片: 網頁原生 H.264 錄影與 HTTP 206 區段串流
    draw_card(
        slide3,
        left=Inches(6.9),
        top=Inches(1.5),
        width=Inches(5.8),
        height=Inches(5.2),
        title="🎬 原生 H.264 與 HTTP 206 區段重播",
        body_list=[
            "• 自動化 OpenH264 部署：系統在 Windows 執行時，若偵測缺少 H.264 編碼，會全自動從 Cisco 官方伺服器下載並極速加載 OpenH264 DLL，實現完全免手動配置的 avc1 本地影片編碼存檔。",
            "• HTTP 206 Range Requests：重播串流路由手動實現 Range 分段傳輸與 206 響應，解決了傳統 HTTP 200 回傳整份影片導致瀏覽器無法快進/倒退或卡死的缺陷，實現了「影片秒開」與「自由拖曳進度條（Seeking）」。"
        ],
        accent_color=CYBER_BLUE
    )

    # ------------------ SLIDE 4: 架構與技術棧 ------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_background(slide4)
    add_header(slide4, "💻 系統架構與技術棧 (Tech Stack)")
    
    # 左卡片: 前端
    draw_card(
        slide4,
        left=Inches(0.6),
        top=Inches(1.5),
        width=Inches(5.8),
        height=Inches(5.2),
        title="💻 網頁控制台 (Frontend UI)",
        body_list=[
            "• 科技美學設計：極簡 Cyberpunk 科技暗色調 UI 面板，極致和諧的視覺搭配與流暢的動態回饋。",
            "• 即時串流播放器：原生 WebSocket 接收影像並通過 HTML5 Image Canvas 即時高效繪製。",
            "• 模組化組件：包含「安全事件歷史日誌（可點擊看截圖與重播）」、「一鍵發送測試警報模組」與「系統設定調整面板」。",
            "• 開發技術：基於原生 HTML5、高性能 CSS3 變量設定、JavaScript 與 Native WebSocket，無笨重框架負擔。"
        ],
        accent_color=CYBER_BLUE
    )
    
    # 右卡片: 後端
    draw_card(
        slide4,
        left=Inches(6.9),
        top=Inches(1.5),
        width=Inches(5.8),
        height=Inches(5.2),
        title="⚙️ 後端引擎與安全資料庫 (Backend Engine)",
        body_list=[
            "• API 核心框架：基於 FastAPI 高性能非同步 Web 框架，提供高併發、極速 JSON API 與 WebSocket 端點服務。",
            "• AI 影像處理：利用 OpenCV (VideoCapture & VideoWriter) 提供基礎影像採集、環形快取、H.264 影片合成；結合 PyTorch YOLOv8 提供人形高精判定。",
            "• 安全警報管道：支援 LINE Notify (截圖與文字)、Discord Webhook (高階 Rich Embed 卡片) 與 Email (MIME SMTP 結合 STARTTLS/SSL)。",
            "• 安全資料庫：採用本機 SQLite 獨立資料庫，所有查詢全面實施 SQL 參數化查詢，徹底根絕 SQL 注入攻擊。"
        ],
        accent_color=CYBER_GREEN
    )

    # ------------------ SLIDE 5: 修改與演進歷程 ------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_background(slide5)
    add_header(slide5, "⏳ 專案修改與演進歷程 (Evolution)")
    
    # 演進歷程採用 4 個 2x2 格局展示階段，以更為清晰精簡地體現四個重要 commit 階段
    w_card = Inches(5.8)
    h_card = Inches(2.4)
    
    # 階段 1
    draw_card(
        slide5,
        left=Inches(0.6),
        top=Inches(1.5),
        width=w_card,
        height=h_card,
        title="Phase 1: 專案起步與目錄架構搭建",
        body_list=[
            "• 確立基本目錄架構，包含 static 前端資源與 storage 本地儲存資料夾。",
            "• 設置完善的 .gitignore，確保 SQLite 資料庫 (nvr.db) 與含有私密 Token 的 config.json 不意外上傳至 GitHub。"
        ],
        accent_color=CYBER_GRAY
    )
    
    # 階段 2
    draw_card(
        slide5,
        left=Inches(6.9),
        top=Inches(1.5),
        width=w_card,
        height=h_card,
        title="Phase 2: 線程解耦與資料庫防丟重構",
        body_list=[
            "• 將 WebSocket send 傳送從局部捕獲解耦，徹底修復瀏覽器斷線時造成無限錯誤 Log 與 CPU 吃滿的缺陷。",
            "• 重構資料庫寫入：由「錄影結束才記錄」改為「偵測入侵當下立即寫入」，保障極致精準的辨識置信度，且能預防異常斷電丟失記錄。"
        ],
        accent_color=CYBER_BLUE
    )
    
    # 階段 3
    draw_card(
        slide5,
        left=Inches(0.6),
        top=Inches(4.3),
        width=w_card,
        height=h_card,
        title="Phase 3: 瀏覽器原生播放與分段快進實裝",
        body_list=[
            "• 將影片編碼改為 avc1 (H.264)，實裝 OpenH264 DLL Windows 自動化極速下載，使網頁原生支援直接播放錄影檔。",
            "• 手動為 uvicorn 實裝 HTTP 206 Range 範圍請求串流路由，徹底解決重播無法拉動進度條與拖曳 Seeking 卡死的頑疾。"
        ],
        accent_color=CYBER_GREEN
    )
    
    # 階段 4
    draw_card(
        slide5,
        left=Inches(6.9),
        top=Inches(4.3),
        width=w_card,
        height=h_card,
        title="Phase 4: 多管道通知警報與精細化錄影防護",
        body_list=[
            "• 實裝可自訂的「單次錄影最大上限時間（max_recording_duration）」，支援自動切分檔案無縫續錄，防範無限錄影塞爆磁碟。",
            "• 完成 LINE Notify、Discord Rich Embed 與 SMTP Email 三管道異步警報，並配有 Cooldown 冷卻計時，防止重疊警報轟炸。"
        ],
        accent_color=CYBER_RED
    )

    # ------------------ SLIDE 6: 遇到的關鍵問題與解決方法 ------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_background(slide6)
    add_header(slide6, "🛠️ 遇到的關鍵工程問題與技術解決方案")
    
    w_3_card = Inches(3.8)
    h_3_card = Inches(5.2)
    
    # 問題 1: 瀏覽器無法播放影片
    draw_card(
        slide6,
        left=Inches(0.6),
        top=Inches(1.5),
        width=w_3_card,
        height=h_3_card,
        title="❌ HTML5 原生拒播 OpenCV 影片",
        body_list=[
            "【痛點成因】",
            "OpenCV 預設的影片格式 (如 XVID、MJPG) 在現代瀏覽器 Chrome/Edge 的原生 <video> 元素中完全無法直接解碼，重播時會顯示檔案損毀或解碼失敗。",
            "【技術解法】",
            "1. 將編碼格式強制指定為 avc1 (H.264) 目標 Codec。",
            "2. 針對 Windows 系統編寫自動引導下載機制，首次運行自動自 Cisco 下載並載入 OpenH264 DLL，實現「零設定零痛苦」的網頁原生錄影解碼。"
        ],
        accent_color=CYBER_RED
    )
    
    # 問題 2: 影片重播無法拖曳進度條
    draw_card(
        slide6,
        left=Inches(4.76),
        top=Inches(1.5),
        width=w_3_card,
        height=h_3_card,
        title="❌ 重播影片無法任意拖曳進度 (Seeking)",
        body_list=[
            "【痛點成因】",
            "使用 FastAPI 的 FileResponse 路由為 HTTP 200，瀏覽器播放影片時必須一口氣下載全部內容，且 HTML5 Player 無法取得隨機區段 (Range)，導致快進、拉動進度條時完全失效卡死。",
            "【技術解法】",
            "1. 手動編寫 HTTP 206 處理程序，解析瀏覽器發送的 HTTP Range 請求標頭 (bytes=start-end)。",
            "2. 設計「二進位分塊生成器 (Chunk Generator)」並傳回 StreamingResponse，加上對應的 Content-Range、Accept-Ranges 標頭，完美解鎖拖曳進度條功能。"
        ],
        accent_color=CYBER_RED
    )
    
    # 問題 3: 斷線引發 Log 狂飆與 CPU 塞滿
    draw_card(
        slide6,
        left=Inches(8.93),
        top=Inches(1.5),
        width=w_3_card,
        height=h_3_card,
        title="❌ 網頁斷線引發線程崩潰與 Log 狂飆",
        body_list=[
            "【痛點成因】",
            "當監控面板被關閉或重載時，WebSocket 連線已斷開，但背景影像擷取發送迴圈無法即時感知斷連，陷入無限迴圈捕捉並嘗試 send，每毫秒大量噴出 Exception 錯誤日誌並把 CPU 負載吃滿。",
            "【技術解法】",
            "1. 將 WebSocket send 與本地擷取邏輯適度解耦。",
            "2. 在最外層精準實裝 try-except WebSocketDisconnect 機制，一旦客戶端斷開就立刻主動關閉擷取線程、安全銷毀寫入器並優雅釋放相機與系統資源。"
        ],
        accent_color=CYBER_RED
    )

    # ------------------ SLIDE 7: 未來擴充想法與展望 ------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_background(slide7)
    add_header(slide7, "💡 未來擴充想法與展望 (Future Discussion)")
    
    # 未來展望 3 張卡片
    draw_card(
        slide7,
        left=Inches(0.6),
        top=Inches(1.5),
        width=w_3_card,
        height=h_3_card,
        title="💡 AI 人臉辨識熟人白名單",
        body_list=[
            "【升級功能】",
            "導入輕量人臉辨識模組 (如 FaceNet 或 InsightFace)。建立親友與熟人白名單資料庫。",
            "【應用價值】",
            "當系統判定有人入侵時，若經人臉特徵比對出是已登記的家庭成員或親人，會自動「抑制 (Suppress) 手機 LINE 與 Discord 警報發送」，杜絕不必要的重複通知騷擾，但背景依舊維持錄影存檔，確保安全防線。"
        ],
        accent_color=CYBER_GREEN
    )
    
    draw_card(
        slide7,
        left=Inches(4.76),
        top=Inches(1.5),
        width=w_3_card,
        height=h_3_card,
        title="💡 雙備份雲端加密同步",
        body_list=[
            "【升級功能】",
            "與主流雲端儲存 API（AWS S3、Google Drive 或 OneDrive）進行非同步異步對接。",
            "【應用價值】",
            "當 AI 判定入侵並完成 5 秒回溯錄影存檔後，立即在背景將該影片加密上傳至雲端。這極大地提升了實體防護安全，即便宵小竊賊將實體的 NVR 主機硬碟整台破壞或抱走，雲端仍留有入侵者最清晰的 5 秒前置特寫畫面。"
        ],
        accent_color=CYBER_BLUE
    )
    
    draw_card(
        slide7,
        left=Inches(8.93),
        top=Inches(1.5),
        width=w_3_card,
        height=h_3_card,
        title="💡 多路 RTSP 與 JWT 安全繫結",
        body_list=[
            "【升級功能】",
            "1. 擴充多路 RTSP/ONVIF 監控頻道，支援網頁多分割畫面同時即時監控。",
            "2. 導入 JWT Token 身分驗證機制與 SSL 安全傳輸協定。",
            "【應用價值】",
            "支援一般居家多門窗、前庭後院的多相機監控需求；加上 JWT 安全認證與傳輸加密，杜絕任何外部未授權用戶或網路駭客對居家私密監視畫面的窺視竊聽。"
        ],
        accent_color=CYBER_GREEN
    )

    # 儲存簡報
    filename = "AI_Soft_NVR_Surveillance_System.pptx"
    prs.save(filename)
    print(f"成功生成簡報！已儲存至: {os.path.abspath(filename)}")

if __name__ == "__main__":
    create_presentation()
